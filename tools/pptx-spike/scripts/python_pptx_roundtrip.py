from __future__ import annotations

import sys
from pathlib import Path


def add_deps() -> None:
    root = Path(__file__).resolve().parents[1]
    for deps in [root / "vendor" / "python", root / ".deps" / "python"]:
        try:
            has_package = (deps / "pptx" / "__init__.py").exists()
        except PermissionError:
            has_package = False
        if has_package:
            sys.path.insert(0, str(deps))
            break


add_deps()

from pptx import Presentation


def roundtrip(source: Path, target: Path) -> dict:
    prs = Presentation(source)
    extraction = {
        "slide_count": len(prs.slides),
        "layouts_visible": len(prs.slide_layouts),
        "masters_visible": len(prs.slide_masters),
        "slides": [],
    }
    for idx, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "index": idx,
            "layout_name": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "shapes": [],
        }
        for shape in slide.shapes:
            info = {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "has_text": bool(getattr(shape, "has_text_frame", False)),
                "has_table": bool(getattr(shape, "has_table", False)),
                "has_chart": bool(getattr(shape, "has_chart", False)),
            }
            if info["has_text"]:
                info["text"] = shape.text
            slide_info["shapes"].append(info)
        extraction["slides"].append(slide_info)

    # Simulate editable regeneration by changing text in existing shapes.
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and "FY2025" in shape.text:
                shape.text = shape.text.replace("FY2025", "FY2026")
            if getattr(shape, "has_text_frame", False) and "Three high-priority" in shape.text:
                shape.text = shape.text.replace("Three high-priority", "Two high-priority")

    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(target)
    return extraction


if __name__ == "__main__":
    import json

    result = roundtrip(Path(sys.argv[1]), Path(sys.argv[2]))
    print(json.dumps(result, indent=2))
