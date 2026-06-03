from __future__ import annotations

import sys
import zipfile
from pathlib import Path


def add_deps() -> None:
    root = Path(__file__).resolve().parents[1]
    for deps in [root / "vendor" / "python", root / ".deps" / "python"]:
        try:
            has_package = (deps / "pptx" / "__init__.py").exists()
        except PermissionError:
            has_package = False
        if has_package:
            sys.path.insert(0, str(deps))
            break


add_deps()

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def create_fixture(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "FY2025 Audit Committee"
    title_slide.placeholders[1].text = "Prior-year template | Confidential"
    banner = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.35))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(31, 78, 121)
    banner.line.color.rgb = RGBColor(31, 78, 121)

    summary_layout = prs.slide_layouts[5]
    summary = prs.slides.add_slide(summary_layout)
    summary.shapes.title.text = "Executive Summary"
    textbox = summary.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(5.75), Inches(4.7))
    tf = textbox.text_frame
    tf.text = "Audit scope completed across core financial controls."
    for text in [
        "Three high-priority findings remain open.",
        "Management remediation cadence improved quarter over quarter.",
        "Committee attention requested for access governance.",
    ]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    rows, cols = 4, 3
    table_shape = summary.shapes.add_table(rows, cols, Inches(7.0), Inches(1.25), Inches(5.3), Inches(2.0))
    table = table_shape.table
    headers = ["Risk", "Open", "Closed"]
    for col, header in enumerate(headers):
        table.cell(0, col).text = header
    for row, values in enumerate([["High", "3", "2"], ["Medium", "7", "5"], ["Low", "4", "8"]], start=1):
        for col, value in enumerate(values):
            table.cell(row, col).text = value

    chart_data = CategoryChartData()
    chart_data.categories = ["High", "Medium", "Low"]
    chart_data.add_series("Open", (3, 7, 4))
    chart_data.add_series("Closed", (2, 5, 8))
    chart = summary.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.05),
        Inches(3.65),
        Inches(5.1),
        Inches(2.5),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    findings = prs.slides.add_slide(prs.slide_layouts[1])
    findings.shapes.title.text = "Finding Detail | Access Governance"
    findings.placeholders[1].text = (
        "Condition: Quarterly access reviews were not consistently evidenced.\n"
        "Risk: Unauthorized access may persist beyond acceptable timeframes.\n"
        "Recommendation: Standardize evidence capture and escalation."
    )

    prs.save(path)
    inject_notes(path, {
        1: "Speaker note: Preserve committee-facing cover language.",
        2: "Speaker note: Discuss remediation trend and high-risk open findings.",
        3: "Speaker note: Confirm management owner and due date.",
    })


def inject_notes(path: Path, notes_by_slide: dict[int, str]) -> None:
    temp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(temp, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                text = data.decode("utf-8")
                if "notesSlide+xml" not in text:
                    text = text.replace(
                        "</Types>",
                        '<Default Extension="vml" ContentType="application/vnd.openxmlformats-officedocument.vmlDrawing"/>'
                        '<Override PartName="/ppt/notesMasters/notesMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>'
                        '<Override PartName="/ppt/notesSlides/notesSlide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                        '<Override PartName="/ppt/notesSlides/notesSlide2.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                        '<Override PartName="/ppt/notesSlides/notesSlide3.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                        "</Types>",
                    )
                    data = text.encode("utf-8")
            for slide_num in notes_by_slide:
                slide_rels = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
                if item.filename == slide_rels:
                    text = data.decode("utf-8")
                    if "notesSlide" not in text:
                        text = text.replace(
                            "</Relationships>",
                            f'<Relationship Id="rIdNotes{slide_num}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" Target="../notesSlides/notesSlide{slide_num}.xml"/></Relationships>',
                        )
                        data = text.encode("utf-8")
            zout.writestr(item, data)

        # Add minimal notes slide parts and relationships. This is enough for Open XML preservation checks.
        for slide_num, note in notes_by_slide.items():
            note_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Notes Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr/>
        <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t>{note}</a:t></a:r></a:p></p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''
            rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="../slides/slide{slide_num}.xml"/>
</Relationships>'''
            zout.writestr(f"ppt/notesSlides/notesSlide{slide_num}.xml", note_xml)
            zout.writestr(f"ppt/notesSlides/_rels/notesSlide{slide_num}.xml.rels", rels_xml)

    temp.replace(path)


if __name__ == "__main__":
    create_fixture(Path(sys.argv[1]))
